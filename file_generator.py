from __future__ import annotations
import csv
import re
import shutil
import subprocess
import wave
import os
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import urlopen
from log_utils import configure_logging, preview_text

logger = configure_logging(__name__)
FONT_CANONICAL_NAMES = {
    'arial': 'Arial',
    'times new roman': 'Times New Roman',
    'calibri': 'Calibri',
    'verdana': 'Verdana',
    'georgia': 'Georgia',
    'garamond': 'Garamond',
    'comic sans ms': 'Comic Sans MS',
    'courier new': 'Courier New',
    'tahoma': 'Tahoma',
    'helvetica': 'Helvetica',
}
ALLOWED_TEXT_STYLES = {'bold', 'italic', 'underline', 'uppercase', 'lowercase', 'title'}
ALLOWED_ALIGNMENTS = {'left', 'center', 'right', 'justify'}
DETAIL_CATEGORY_ALIASES = {
    'image': 'images',
    'images': 'images',
    'picture': 'images',
    'photo': 'images',
    'table': 'tables',
    'tables': 'tables',
    'chart': 'charts',
    'charts': 'charts',
    'graph': 'graphs',
    'graphs': 'graphs',
    'plot': 'graphs',
    'hyperlink': 'hyperlinks',
    'hyperlinks': 'hyperlinks',
    'link': 'hyperlinks',
    'links': 'hyperlinks',
    'url': 'hyperlinks',
    'urls': 'hyperlinks',
    'header': 'headers',
    'headers': 'headers',
    'footer': 'footers',
    'footers': 'footers',
    'section': 'sections',
    'sections': 'sections',
    'page_number': 'page_numbers',
    'page_numbers': 'page_numbers',
    'table_of_contents': 'tables_of_contents',
    'tables_of_contents': 'tables_of_contents',
    'toc': 'tables_of_contents',
    'index': 'indexes',
    'indexes': 'indexes',
    'bibliography': 'bibliographies',
    'bibliographies': 'bibliographies',
    'citation': 'citations',
    'citations': 'citations',
    'footnote': 'footnotes',
    'footnotes': 'footnotes',
    'note': 'notes',
    'notes': 'notes',
    'margin': 'margins',
    'margins': 'margins',
    'padding': 'paddings',
    'paddings': 'paddings',
    'border': 'borders',
    'borders': 'borders',
    'background': 'backgrounds',
    'backgrounds': 'backgrounds',
    'layout': 'layouts',
    'layouts': 'layouts',
    'template': 'templates',
    'templates': 'templates',
    'theme': 'themes',
    'themes': 'themes',
    'font': 'fonts',
    'fonts': 'fonts',
    'color': 'colors',
    'colors': 'colors',
    'size': 'sizes',
    'sizes': 'sizes',
    'style': 'styles',
    'styles': 'styles',
    'alignment': 'alignments',
    'alignments': 'alignments',
}
IMAGE_FILE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'}
TEXT_FILE_EXTENSIONS = {
    'txt', 'text', 'docx',
    'md', 'markdown', 'html', 'htm', 'css',
    'py', 'js', 'java', 'c', 'cpp', 'rb', 'go', 'rs', 'swift', 'kt', 'ts', 'php',
    'json', 'xml', 'yml', 'yaml', 'toml', 'ini',
    'sh', 'bat', 'ps1',
    'svg',
}
CODE_FILE_EXTENSIONS = {'py', 'js', 'java', 'c', 'cpp', 'rb', 'go', 'rs', 'swift', 'kt', 'ts', 'php'}
CSV_FILE_EXTENSIONS = {'csv'}
PDF_FILE_EXTENSIONS = {'pdf'}
PPT_FILE_EXTENSIONS = {'ppt', 'pptx'}
AUDIO_FILE_EXTENSIONS = {'mp3', 'wav', 'flac', 'aac', 'ogg', 'm4a'}
VIDEO_FILE_EXTENSIONS = {'mp4', 'avi', 'mkv', 'mov', 'wmv', 'webm'}
CLOUD_PROVIDER_ALIASES = {
    'google_drive': 'google_drive',
    'googledrive': 'google_drive',
    'gdrive': 'google_drive',
    'drive': 'google_drive',
    'dropbox': 'dropbox',
    'onedrive': 'onedrive',
    'icloud': 'icloud',
    'box': 'box',
    's3': 's3',
    'amazon_s3': 's3',
    'pcloud': 'pcloud',
    'mega': 'mega',
    'sync': 'sync',
    'sync_com': 'sync',
    'nextcloud': 'nextcloud',
    'owncloud': 'owncloud',
    'tresorit': 'tresorit',
}
RCLONE_REMOTE_DEFAULTS = {
    'google_drive': 'google_drive',
    'dropbox': 'dropbox',
    'onedrive': 'onedrive',
    'icloud': 'icloud',
    'box': 'box',
    's3': 's3',
    'pcloud': 'pcloud',
    'mega': 'mega',
    'sync': 'sync',
    'nextcloud': 'nextcloud',
    'owncloud': 'owncloud',
    'tresorit': 'tresorit',
}
_RCLONE_AVAILABLE: bool | None = None

try:
    from reportlab.pdfgen import canvas as reportlab_canvas
except ImportError:
    reportlab_canvas = None
    logger.warning('reportlab not installed; .pdf write unavailable')

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None
    logger.warning('PyPDF2 not installed; .pdf read unavailable')

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logger.warning('python-pptx not installed; .pptx support unavailable')

try:
    import pandas as pd
except ImportError:
    pd = None
    logger.warning('pandas not installed; .xlsx generation unavailable')

try:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.opc.constants import RELATIONSHIP_TYPE
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.shared import Pt, RGBColor, Inches, Cm, Mm
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    from docx.enum.section import WD_ORIENTATION
except ImportError:
    Document = None
    CT_Tbl = None
    CT_P = None
    OxmlElement = None
    qn = None
    RELATIONSHIP_TYPE = None
    Table = None
    Paragraph = None
    Pt = None
    RGBColor = None
    Inches = None
    Cm = None
    Mm = None
    WD_PARAGRAPH_ALIGNMENT = None
    WD_ORIENTATION = None
    logger.warning('python-docx not installed; .docx generation unavailable')

from intel import generate_image
try:
    import matplotlib.pyplot as plt
except ImportError:
    plt = None
    logger.warning('matplotlib not installed; chart generation unavailable')

## Main agent function that routes to specific file handling functions based on file type and action

def agent(
    text: str,
    file_name: str,
    action: str,
    style: str = '',
    chart_type: str = None,
    chart_data: dict = None,
    format_options: dict | None = None,
    details: object = None,
) -> str:
    original_file_name = file_name
    file_name = _resolve_storage_target(file_name)
    cloud_ref = _parse_cloud_reference(original_file_name)
    use_rclone = cloud_ref is not None and _rclone_is_available()

    if use_rclone and action == 'D' and _is_full_delete_request(text):
        deleted, error_message = _rclone_delete_remote(cloud_ref)
        if deleted:
            local_path = Path(file_name)
            if local_path.exists():
                try:
                    local_path.unlink()
                except Exception:
                    pass
            return f'Deleted file "{original_file_name}" successfully.'
        return f'Error deleting file: {error_message}'

    if use_rclone and action in ('R', 'A'):
        downloaded, error_message = _rclone_download_remote(cloud_ref, file_name)
        if action == 'R' and not downloaded:
            return f'Error reading file: {error_message}'

    if use_rclone and action == 'D' and not _is_full_delete_request(text):
        downloaded, error_message = _rclone_download_remote(cloud_ref, file_name)
        if not downloaded:
            return f'Error deleting content: {error_message}'

    result = _agent_local(
        text=text,
        file_name=file_name,
        action=action,
        style=style,
        chart_type=chart_type,
        chart_data=chart_data,
        format_options=format_options,
        details=details,
    )

    if not use_rclone:
        return result

    if _is_error_result(result):
        return result

    if action in ('W', 'A'):
        upload_file = file_name
        upload_ref = cloud_ref
        if isinstance(result, str):
            candidate = Path(result.strip())
            if candidate.exists():
                upload_file = str(candidate)
                upload_ref = _cloud_ref_with_result_suffix(cloud_ref, candidate)
        uploaded, error_message = _rclone_upload_remote(upload_ref, upload_file)
        if not uploaded:
            return f'Error syncing cloud file: {error_message}'
        if isinstance(result, str):
            return _cloud_display_name(original_file_name, upload_ref)
        return str(original_file_name)

    if action == 'D' and not _is_full_delete_request(text):
        uploaded, error_message = _rclone_upload_remote(cloud_ref, file_name)
        if not uploaded:
            return f'Error syncing cloud file: {error_message}'
        if isinstance(result, str):
            return result + f' Synced to {original_file_name}.'

    return result

def _agent_local(
    text: str,
    file_name: str,
    action: str,
    style: str = '',
    chart_type: str = None,
    chart_data: dict = None,
    format_options: dict | None = None,
    details: object = None,
) -> str:
    # Divides to functions based on file extension and action
    ext = Path(file_name).suffix.lower().lstrip('.')
    options = _normalize_format_options(format_options)
    logger.info(
        'Agent dispatch start | action=%s resolved_file=%s ext=%s text_len=%d text_preview=%s',
        action,
        file_name,
        ext,
        len(str(text)),
        preview_text(text),
    )
    if chart_type and chart_data:
        logger.debug(
            'Agent routed to chart handler | action=%s file=%s chart_type=%s points=%d',
            action,
            file_name,
            chart_type,
            len(chart_data) if isinstance(chart_data, dict) else 0,
        )
        if action in ('W', 'A'):
            output_path = _resolve_chart_output_path(file_name)
            try:
                generate_chart(chart_type, chart_data, str(output_path))
            except Exception as exc:
                logger.exception('Chart generation failed | file=%s chart_type=%s', file_name, chart_type)
                return f'Error generating chart: {exc}'
            return str(output_path)
        if action == 'D':
            return delete(file_name, text)
        return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'

    if ext in TEXT_FILE_EXTENSIONS or ext == '':
        logger.debug('Agent routed to text/docx handler | action=%s file=%s', action, file_name)
        if action == 'W':
            return txt_write(file_name, text, style, options, details=details)
        elif action == 'R':
            a = txt_read(file_name)
            return a
        elif action == 'A':
            return txt_append(file_name, text, style, options, details=details)
        elif action == 'D':
            return delete(file_name, text)
        else:
            logger.error('Invalid action specified: %s', action)
            return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'
    elif ext in ('xlsx', 'xls'):
        logger.debug('Agent routed to xlsx handler | action=%s file=%s', action, file_name)
        if action == 'W':
            return xlsx_write(file_name, text)
        elif action == 'R':
            return xlsx_read(file_name)
        elif action == 'A':
            return xlsx_append(file_name, text)
        elif action == 'D':
            return delete(file_name, text)
        else:
            logger.error('Invalid action specified: %s', action)
            return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'
    elif ext in CSV_FILE_EXTENSIONS:
        logger.debug('Agent routed to csv handler | action=%s file=%s', action, file_name)
        if action == 'W':
            return csv_write(file_name, text)
        elif action == 'R':
            return csv_read(file_name)
        elif action == 'A':
            return csv_append(file_name, text)
        elif action == 'D':
            return delete(file_name, text)
        else:
            logger.error('Invalid action specified: %s', action)
            return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'
    elif ext in PDF_FILE_EXTENSIONS:
        logger.debug('Agent routed to pdf handler | action=%s file=%s', action, file_name)
        if action == 'W':
            return pdf_write(file_name, text)
        elif action == 'R':
            return pdf_read(file_name)
        elif action == 'A':
            return pdf_append(file_name, text)
        elif action == 'D':
            return delete(file_name, text)
        else:
            logger.error('Invalid action specified: %s', action)
            return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'
    elif ext in PPT_FILE_EXTENSIONS:
        logger.debug('Agent routed to ppt handler | action=%s file=%s', action, file_name)
        if action == 'W':
            return ppt_write(file_name, text)
        elif action == 'R':
            return ppt_read(file_name)
        elif action == 'A':
            return ppt_append(file_name, text)
        elif action == 'D':
            return delete(file_name, text)
        else:
            logger.error('Invalid action specified: %s', action)
            return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'
    elif ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp'):
        logger.debug('Agent routed to image handler | action=%s file=%s', action, file_name)
        if action == 'W':
            return image_creation(text, file_name, action)
        elif action == 'R':
            #
            return image_read(text, file_name, action)
        elif action == 'A':
            return image_append(text, file_name, action)
        elif action == 'D':
            return delete(file_name, text)
        else:
            logger.error('Invalid action specified: %s', action)
            return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'
    elif ext in AUDIO_FILE_EXTENSIONS:
        logger.debug('Agent routed to audio handler | action=%s file=%s', action, file_name)
        if action == 'R':
            return media_read(file_name, media_type='audio')
        if action == 'W':
            return media_write(file_name, text, media_type='audio', mode='write')
        if action == 'A':
            return media_write(file_name, text, media_type='audio', mode='append')
        if action == 'D':
            return delete(file_name, text)
        logger.error('Invalid action specified: %s', action)
        return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'
    elif ext in VIDEO_FILE_EXTENSIONS:
        logger.debug('Agent routed to video handler | action=%s file=%s', action, file_name)
        if action == 'R':
            return media_read(file_name, media_type='video')
        if action == 'W':
            return media_write(file_name, text, media_type='video', mode='write')
        if action == 'A':
            return media_write(file_name, text, media_type='video', mode='append')
        if action == 'D':
            return delete(file_name, text)
        logger.error('Invalid action specified: %s', action)
        return 'Error: Invalid action. Use W (write), A (append), R (read), or D (delete).'

    ## here goes graph logic
    else:
        logger.error('Unsupported file type specified: %s', ext)
    return 'Error: Unsupported file type. Please choose a supported extension.'

###--BEGINNING OF TEXT FILE GENERATOR--##

def read(file_name: str) -> str:
    # Helper for external imports to read files based on extension
    ext = Path(file_name).suffix.lower().lstrip('.')
    logger.info('Read helper invoked | file=%s ext=%s', file_name, ext)
    if ext in ('xlsx', 'xls'):
        return xlsx_read(file_name)
    if ext in CSV_FILE_EXTENSIONS:
        return csv_read(file_name)
    if ext in PDF_FILE_EXTENSIONS:
        return pdf_read(file_name)
    if ext in PPT_FILE_EXTENSIONS:
        return ppt_read(file_name)
    if ext in AUDIO_FILE_EXTENSIONS:
        return media_read(file_name, media_type='audio')
    if ext in VIDEO_FILE_EXTENSIONS:
        return media_read(file_name, media_type='video')
    return txt_read(file_name)

def txt_write(
    file_name: str,
    text: str,
    style: str = '',
    format_options: dict | None = None,
    details: object = None,
) -> str:
    path = Path(file_name)
    options = _normalize_format_options(format_options)
    logger.info(
        'Text write requested | file=%s ext=%s text_len=%d text_preview=%s',
        file_name,
        path.suffix.lower(),
        len(str(text)),
        preview_text(text),
    )
    try:
        _ensure_parent_dir(path)
    except Exception as exc:
        logger.exception('Failed to prepare directory for %s', file_name)
        return f'Error preparing directory: {exc}'

    if file_name.lower().endswith('.docx'):
        if Document is None:
            msg = 'python-docx is required to generate .docx files. Install with "pip install python-docx".'
            print(msg)
            logger.error(msg)
            return 'Missing dependency: python-docx'
        try:
            doc = Document()
            if style:
                doc.add_paragraph(f"Font Style: {style}", style='Heading 1')
            _append_blocks_to_docx(doc, text, options, details=details, file_name=file_name)
            doc.save(file_name)
            logger.info('Created docx file %s | tables=%d paragraphs=%d', file_name, len(doc.tables), len(doc.paragraphs))
            return file_name
        except Exception as exc:
            logger.exception('Failed to create docx file %s', file_name)
            return f'Error creating docx: {exc}'
    try:
        with open(file_name, 'w', encoding='utf-8') as file:
            text_with_details = _append_txt_details(text, details)
            styled_text = _apply_txt_style(text_with_details, style, options)
            file.write(styled_text)
            logger.info('Created text file %s | bytes=%d', file_name, len(styled_text.encode("utf-8", errors="replace")))
            return file_name
    except Exception as exc:
        logger.exception('Failed to write text file %s', file_name)
        return f'Error writing text file: {exc}'


def txt_read(file_name: str) -> str:
    path = Path(file_name)
    logger.info('Text read requested | file=%s ext=%s', file_name, path.suffix.lower())
    try:
        if path.suffix.lower() == '.docx':
            if Document is None:
                msg = 'python-docx is required to read .docx files. Install with "pip install python-docx".'
                print(msg)
                logger.error(msg)
                return 'Missing dependency: python-docx'
            doc = Document(file_name)
            context = _read_docx_content(doc)
            logger.info(
                'Read docx file %s | tables=%d paragraphs=%d extracted_len=%d',
                file_name,
                len(doc.tables),
                len(doc.paragraphs),
                len(context),
            )
            return context


        with open(file_name, 'r', encoding='utf-8', errors='replace') as f:
            context = f.read()
        logger.info('Read text file %s | chars=%d preview=%s', file_name, len(context), preview_text(context))
        return context
    except Exception as exc:
        logger.exception('Failed to read file %s', file_name)
        return f'Error reading file: {exc}'

def txt_append(
    file_name: str,
    text: str,
    style: str = '',
    format_options: dict | None = None,
    details: object = None,
) -> str:
    path = Path(file_name)
    options = _normalize_format_options(format_options)
    logger.info(
        'Text append requested | file=%s ext=%s text_len=%d text_preview=%s',
        file_name,
        path.suffix.lower(),
        len(str(text)),
        preview_text(text),
    )
    try:
        _ensure_parent_dir(path)
    except Exception as exc:
        logger.exception('Failed to prepare directory for %s', file_name)
        return f'Error preparing directory: {exc}'

    if path.suffix.lower() == '.docx':
        if Document is None:
            msg = 'python-docx is required to generate .docx files. Install with "pip install python-docx".'
            print(msg)
            logger.error(msg)
            return 'Missing dependency: python-docx'
        try:
            doc = Document(file_name) if path.exists() else Document()
            _append_blocks_to_docx(doc, text, options, details=details, file_name=file_name)
            doc.save(file_name)
            logger.info('Appended to docx file %s | tables=%d paragraphs=%d', file_name, len(doc.tables), len(doc.paragraphs))
            return file_name
        except Exception as exc:
            logger.exception('Failed to append to docx file %s', file_name)
            return f'Error appending to docx: {exc}'

    try:
        prefix = '\n' if path.exists() and path.stat().st_size > 0 else ''
        with open(file_name, 'a', encoding='utf-8') as file:
            text_with_details = _append_txt_details(text, details)
            styled_text = _apply_txt_style(text_with_details, style, options)
            file.write(prefix + styled_text)
            logger.info(
                'Appended to text file %s | prefix_added=%s appended_chars=%d',
                file_name,
                bool(prefix),
                len(styled_text),
            )
            return file_name
    except FileNotFoundError:
        try:
            with open(file_name, 'w', encoding='utf-8') as f:
                text_with_details = _append_txt_details(text, details)
                styled_text = _apply_txt_style(text_with_details, style, options)
                f.write(styled_text)
            logger.info('Append target missing; created new text file %s | chars=%d', file_name, len(styled_text))
            return file_name
        except Exception as exc:
            logger.exception('Failed to create new file %s after FileNotFoundError', file_name)
            return f'Error creating new file after FileNotFoundError: {exc}'
    except Exception as exc2:
        logger.exception('Failed to append to file %s', file_name)
        return f'Error appending to file: {exc2}'

def _apply_txt_style(text: str, style: str, format_options: dict | None = None) -> str:
    options = _normalize_format_options(format_options)
    selected = (style or '').strip().lower()
    output = text
    if selected in ('arial', 'calibri'):
        output = text
    elif selected in ('times new roman', 'georgia'):
        output = text.title()
    elif selected == 'verdana':
        output = text.upper()
    elif selected == 'garamond':
        output = text.lower()
    elif selected == 'comic sans ms':
        output = _comic_case(text)

    output = _apply_text_case_styles(output, options.get('styles'))
    output = _apply_txt_alignment(output, options.get('alignment', ''))
    return output

def _comic_case(text: str) -> str:
    chars: list[str] = []
    upper = True
    for char in text:
        if char.isalpha():
            chars.append(char.upper() if upper else char.lower())
            upper = not upper
        else:
            chars.append(char)
    return ''.join(chars)

def _apply_txt_alignment(text: str, alignment: str, width: int | None = None) -> str:
    mode = str(alignment or '').strip().lower()
    if mode not in ALLOWED_ALIGNMENTS:
        return text

    lines = text.splitlines()
    if width is None:
        non_empty_lengths = [len(line.strip()) for line in lines if line.strip()]
        longest = max(non_empty_lengths, default=0)
        width = min(max(longest + 4, 40), 120)
    if width <= 0:
        return text

    aligned_lines: list[str] = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            aligned_lines.append('')
            continue
        if len(stripped) >= width:
            aligned_lines.append(stripped)
            continue
        if mode == 'center':
            aligned_lines.append(stripped.center(width))
        elif mode == 'right':
            aligned_lines.append(stripped.rjust(width))
        elif mode == 'justify':
            words = stripped.split()
            if len(words) <= 1:
                aligned_lines.append(stripped)
            else:
                total_chars = sum(len(word) for word in words)
                spaces_needed = max(width - total_chars, len(words) - 1)
                gaps = len(words) - 1
                if gaps <= 0:
                    aligned_lines.append(stripped)
                    continue
                min_spaces, extra = divmod(spaces_needed, gaps)
                pieces: list[str] = []
                for idx, word in enumerate(words[:-1]):
                    pad = min_spaces + (1 if idx < extra else 0)
                    pieces.append(word + (' ' * pad))
                pieces.append(words[-1])
                aligned_lines.append(''.join(pieces).rstrip())
        else:
            aligned_lines.append(stripped)
    return '\n'.join(aligned_lines)

def _apply_text_case_styles(text: str, styles: object) -> str:
    if not styles:
        return text
    if not isinstance(styles, list):
        return text

    output = text
    normalized = [str(item).lower() for item in styles]
    if 'uppercase' in normalized:
        output = output.upper()
    elif 'lowercase' in normalized:
        output = output.lower()
    elif 'title' in normalized:
        output = output.title()
    return output

def _normalize_format_options(format_options: dict | None) -> dict[str, object]:
    if not isinstance(format_options, dict):
        return {}

    normalized: dict[str, object] = {}

    raw_font = str(format_options.get('font', '')).strip()
    if raw_font:
        normalized_font = FONT_CANONICAL_NAMES.get(raw_font.lower(), raw_font)
        normalized['font'] = normalized_font[:80]

    raw_color = format_options.get('color')
    if _parse_color_to_rgb(raw_color) is not None:
        normalized['color'] = str(raw_color).strip().lower()

    raw_size = format_options.get('size')
    if raw_size is not None:
        try:
            size = float(raw_size)
        except (TypeError, ValueError):
            size = None
        if size is not None and 6 <= size <= 96:
            normalized['size'] = size

    raw_styles = format_options.get('styles')
    style_items: list[str] = []
    if isinstance(raw_styles, str):
        style_items = [item.strip().lower() for item in raw_styles.split(',') if item.strip()]
    elif isinstance(raw_styles, (list, tuple, set)):
        style_items = [str(item).strip().lower() for item in raw_styles if str(item).strip()]
    filtered_styles = [item for item in style_items if item in ALLOWED_TEXT_STYLES]
    if filtered_styles:
        ordered_styles = list(dict.fromkeys(filtered_styles))
        case_styles = [style for style in ordered_styles if style in ('uppercase', 'lowercase', 'title')]
        if len(case_styles) > 1:
            first_case = case_styles[0]
            ordered_styles = [style for style in ordered_styles if style not in ('uppercase', 'lowercase', 'title')]
            ordered_styles.append(first_case)
        normalized['styles'] = ordered_styles

    raw_alignment = str(format_options.get('alignment', '')).strip().lower()
    if raw_alignment in ALLOWED_ALIGNMENTS:
        normalized['alignment'] = raw_alignment

    return normalized

def delete(file_name: str, delete_request: str = '') -> str:
    request = (delete_request or '').strip()
    lower_request = request.lower()
    delete_file_aliases = {'file', 'entire file', 'full file', 'whole file', 'delete file', 'remove file'}
    logger.info(
        'Delete requested | file=%s request_len=%d request_preview=%s',
        file_name,
        len(request),
        preview_text(request),
    )

    if not request or lower_request in delete_file_aliases:
        try:
            p = Path(file_name)
            p.unlink()
            logger.info('Deleted file %s', file_name)
            return f'Deleted file "{file_name}" successfully.'
        except Exception as exc:
            logger.exception('Failed to delete file %s', file_name)
            return f'Error deleting file: {exc}'

    path = Path(file_name)
    ext = path.suffix.lower()
    if ext == '.docx':
        return _delete_docx_content(file_name, request)
    if ext in {'', '.txt'} or ext.lstrip('.') in TEXT_FILE_EXTENSIONS or ext.lstrip('.') in CSV_FILE_EXTENSIONS:
        return _delete_text_content(file_name, request)

    logger.warning('Content delete requested for unsupported file type %s: %s', ext, file_name)
    return f'Error deleting content: targeted content delete is not supported for "{path.suffix or "unknown"}" files.'

def _append_blocks_to_docx(
    doc,
    text: str,
    format_options: dict | None = None,
    details: object = None,
    file_name: str = '',
) -> None:
    options = format_options or {}
    blocks = _parse_docx_blocks(str(text))
    logger.debug('DOCX block parser returned %d block(s)', len(blocks))
    for block in blocks:
        if block['type'] == 'paragraph':
            paragraph_value = _apply_text_case_styles(str(block['value']), options.get('styles'))
            paragraph = doc.add_paragraph(paragraph_value)
            _apply_docx_formatting(paragraph, options)
            logger.debug('Added paragraph block to docx | preview=%s', preview_text(block['value']))
            continue
        rows = block['rows']
        if not isinstance(rows, list) or not rows:
            logger.debug('Skipped empty or invalid table block during docx append')
            continue
        _add_table_to_docx(doc, rows, options)

    _append_docx_details(doc, details, options, file_name)

def _append_txt_details(text: str, details: object) -> str:
    base = str(text or '')
    detail_items = list(_iter_detail_items(details))
    if not detail_items:
        return base

    lines = [base] if base else []
    lines.append('')
    lines.append('Additional details:')
    for category, value in detail_items:
        lines.append(f'- {category}: {value}')
    return '\n'.join(lines).strip('\n')

def _iter_detail_items(details: object):
    if not details:
        return

    if isinstance(details, dict):
        items_value = details.get('items')
        if isinstance(items_value, list):
            for item in items_value:
                yield from _iter_detail_items([item])
            return
        for key, value in details.items():
            if key == 'items':
                continue
            if isinstance(value, (list, tuple, set)):
                for item_value in value:
                    category = _normalize_detail_category(str(key))
                    normalized_value = str(item_value or '').strip()
                    if category and normalized_value:
                        yield category, normalized_value
            else:
                category = _normalize_detail_category(str(key))
                normalized_value = str(value or '').strip()
                if category and normalized_value:
                    yield category, normalized_value
        return

    if isinstance(details, (list, tuple, set)):
        for item in details:
            category = ''
            value = ''
            if isinstance(item, dict):
                category = _normalize_detail_category(str(item.get('category', '')))
                value = str(item.get('value', '')).strip()
            elif isinstance(item, (list, tuple)) and len(item) >= 2:
                category = _normalize_detail_category(str(item[0]))
                value = str(item[1]).strip()
            if category and value:
                yield category, value

def _normalize_detail_category(raw: str) -> str:
    key = re.sub(r'[\s\-]+', '_', str(raw or '').strip().lower())
    return DETAIL_CATEGORY_ALIASES.get(key, '')

def _append_docx_details(doc, details: object, format_options: dict | None, file_name: str) -> None:
    options = format_options or {}
    image_count = 0
    chart_count = 0

    for category, value in _iter_detail_items(details):
        if category in {'fonts', 'colors', 'sizes', 'styles', 'alignments'}:
            continue

        if category == 'tables':
            rows = _parse_table_rows_from_detail(value)
            if rows:
                _add_table_to_docx(doc, rows, options)
            else:
                _append_docx_detail_note(doc, category, value, options)
            continue

        if category == 'images':
            image_count += 1
            image_path = _resolve_detail_image_path(value, file_name, image_count)
            if image_path is None:
                _append_docx_detail_note(doc, category, value, options)
            else:
                _insert_docx_image(doc, image_path, options)
            continue

        if category in {'charts', 'graphs'}:
            chart_count += 1
            chart_path = _resolve_chart_detail_image(value, file_name, chart_count)
            if chart_path is None:
                _append_docx_detail_note(doc, category, value, options)
            else:
                _insert_docx_image(doc, chart_path, options)
            continue

        if category == 'hyperlinks':
            text, url = _parse_hyperlink_detail(value)
            if url:
                paragraph = doc.add_paragraph()
                _add_docx_hyperlink(paragraph, url, text or url)
                _apply_docx_formatting(paragraph, options)
            else:
                _append_docx_detail_note(doc, category, value, options)
            continue

        if category == 'headers':
            _set_docx_header_text(doc, value, options)
            continue

        if category == 'footers':
            _set_docx_footer_text(doc, value, options)
            continue

        if category == 'page_numbers':
            _add_docx_page_numbers(doc, options)
            continue

        if category == 'tables_of_contents':
            _add_docx_toc_field(doc, options)
            continue

        if category == 'sections':
            doc.add_page_break()
            paragraph = doc.add_paragraph(value)
            _apply_docx_formatting(paragraph, options)
            continue

        if category == 'margins':
            _apply_docx_margins(doc, value)
            continue

        if category == 'paddings':
            _apply_docx_paddings(doc, value)
            continue

        if category == 'layouts':
            _apply_docx_layout(doc, value)
            continue

        if category == 'templates':
            _apply_docx_template(doc, value, options)
            continue

        if category == 'themes':
            _apply_docx_theme(doc, value, options)
            continue

        if category == 'backgrounds':
            _apply_docx_background(doc, value)
            continue

        if category == 'borders':
            _apply_docx_borders(doc, value, options)
            continue

        if category == 'indexes':
            _add_docx_index_block(doc, value, options)
            continue

        if category == 'bibliographies':
            _add_docx_reference_block(doc, 'Bibliography', value, options, numbered=False)
            continue

        if category == 'citations':
            _add_docx_reference_block(doc, 'Citations', value, options, numbered=True)
            continue

        if category == 'footnotes':
            _add_docx_reference_block(doc, 'Footnotes', value, options, numbered=True)
            continue

        _append_docx_detail_note(doc, category, value, options)

def _append_docx_detail_note(doc, category: str, value: str, format_options: dict | None) -> None:
    label = category.replace('_', ' ').title()
    paragraph = doc.add_paragraph(f'{label}: {value}')
    _apply_docx_formatting(paragraph, format_options)

def _add_table_to_docx(doc, rows: list[list[str]], format_options: dict | None) -> None:
    if not rows:
        return
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    try:
        table.style = 'Table Grid'
    except Exception:
        pass  # Keep default style if Table Grid is unavailable
    logger.debug('Added table block to docx | rows=%d cols=%d', len(rows), len(rows[0]))
    for row_index, row in enumerate(rows):
        if not isinstance(row, list):
            continue
        for col_index, value in enumerate(row):
            transformed = _apply_text_case_styles(str(value), (format_options or {}).get('styles'))
            cell = table.cell(row_index, col_index)
            cell.text = transformed
            for paragraph in cell.paragraphs:
                _apply_docx_formatting(paragraph, format_options)

def _parse_table_rows_from_detail(value: str) -> list[list[str]]:
    raw = str(value or '').strip()
    if not raw:
        return []

    row_chunks = [part.strip() for part in raw.splitlines() if part.strip()]
    if len(row_chunks) <= 1:
        row_chunks = [part.strip() for part in raw.split('/') if part.strip()]
    if len(row_chunks) <= 1 and ';' in raw:
        row_chunks = [part.strip() for part in raw.split(';') if part.strip()]
    if not row_chunks:
        return []

    rows: list[list[str]] = []
    for row_chunk in row_chunks:
        cells = _split_explicit_table_row(row_chunk)
        if len(cells) < 2:
            return []
        rows.append(cells)

    if len(rows) > 1 and _is_markdown_separator_row(rows[1]):
        rows = [rows[0]] + rows[2:]
    if not rows:
        return []

    max_cols = max((len(row) for row in rows), default=0)
    if max_cols < 2:
        return []
    return [row + [''] * (max_cols - len(row)) for row in rows]

def _split_explicit_table_row(row_text: str) -> list[str]:
    candidate = row_text.strip()
    if '|' in candidate:
        return [part.strip() for part in candidate.strip('|').split('|')]
    if '\t' in candidate:
        return [part.strip() for part in candidate.split('\t')]
    if ',' in candidate:
        try:
            parsed = next(csv.reader([candidate], delimiter=','))
            return [part.strip() for part in parsed]
        except Exception:
            return []
    if ';' in candidate:
        try:
            parsed = next(csv.reader([candidate], delimiter=';'))
            return [part.strip() for part in parsed]
        except Exception:
            return []
    return []

def _resolve_detail_image_path(value: str, file_name: str, index: int) -> Path | None:
    raw = str(value or '').strip()
    if not raw:
        return None

    lower_raw = raw.lower()
    if lower_raw.startswith('path:'):
        local_path = raw.split(':', 1)[1].strip()
        return _resolve_local_image_path(local_path)

    if lower_raw.startswith('http://') or lower_raw.startswith('https://'):
        return _download_image_detail(raw, file_name, index)

    if lower_raw.startswith('prompt:'):
        prompt = raw.split(':', 1)[1].strip()
        return _generate_image_from_prompt(prompt, file_name, index)

    local = _resolve_local_image_path(raw)
    if local is not None:
        return local

    return _generate_image_from_prompt(raw, file_name, index)

def _resolve_local_image_path(value: str) -> Path | None:
    candidate = Path(str(value or '').strip()).expanduser()
    if not candidate.is_absolute():
        candidate = (Path.cwd() / candidate).resolve()
    if candidate.exists() and candidate.suffix.lower() in IMAGE_FILE_EXTENSIONS:
        return candidate
    return None

def _download_image_detail(url: str, file_name: str, index: int) -> Path | None:
    parsed = urlparse(url)
    suffix = Path(parsed.path).suffix.lower()
    if suffix not in IMAGE_FILE_EXTENSIONS:
        suffix = '.png'
    target = _build_detail_asset_path(file_name, f'detail_image_{index}', suffix)
    try:
        with urlopen(url, timeout=20) as response:
            payload = response.read()
        if not payload:
            return None
        target.write_bytes(payload)
        return target if target.exists() else None
    except Exception:
        logger.warning('Failed to download detail image from URL: %s', url)
        return None

def _generate_image_from_prompt(prompt: str, file_name: str, index: int) -> Path | None:
    normalized_prompt = str(prompt or '').strip()
    if not normalized_prompt:
        return None
    target = _build_detail_asset_path(file_name, f'detail_image_{index}', '.png')
    result = generate_image(normalized_prompt, str(target), 'W')
    if isinstance(result, str):
        generated = Path(result.strip()).expanduser()
        if not generated.is_absolute():
            generated = (Path.cwd() / generated).resolve()
        if generated.exists() and generated.suffix.lower() in IMAGE_FILE_EXTENSIONS:
            return generated
    if target.exists():
        return target
    return None

def _build_detail_asset_path(file_name: str, stem_suffix: str, extension: str) -> Path:
    base = Path(file_name) if file_name else Path('output.docx')
    if not base.parent.exists():
        base.parent.mkdir(parents=True, exist_ok=True)
    safe_extension = extension if extension.startswith('.') else f'.{extension}'
    return base.parent / f'{base.stem}_{stem_suffix}{safe_extension}'

def _resolve_chart_detail_image(value: str, file_name: str, index: int) -> Path | None:
    raw = str(value or '').strip()
    if not raw:
        return None

    lower_raw = raw.lower()
    if lower_raw.startswith('path:'):
        local = _resolve_local_image_path(raw.split(':', 1)[1].strip())
        if local is not None:
            return local
    elif lower_raw.startswith('http://') or lower_raw.startswith('https://'):
        downloaded = _download_image_detail(raw, file_name, index)
        if downloaded is not None:
            return downloaded
    else:
        local = _resolve_local_image_path(raw)
        if local is not None:
            return local

    parsed = _parse_chart_detail_spec(raw)
    if parsed is None:
        return None

    chart_type, chart_data = parsed
    output_path = _build_detail_asset_path(file_name, f'detail_chart_{index}', '.png')
    try:
        generate_chart(chart_type, chart_data, str(output_path))
    except Exception:
        logger.warning('Failed to generate chart from detail spec: %s', preview_text(value))
        return None
    return output_path if output_path.exists() else None

def _parse_chart_detail_spec(value: str) -> tuple[str, dict[str, float]] | None:
    text = str(value or '').strip()
    if not text:
        return None

    match = re.match(r'^(line|bar|pie|scatter)\s*[\|:]\s*(.+)$', text, flags=re.IGNORECASE)
    if match:
        chart_type = match.group(1).lower()
        data = _parse_chart_data_pairs(match.group(2))
        return (chart_type, data) if data else None

    type_match = re.search(r'type\s*=\s*(line|bar|pie|scatter)', text, flags=re.IGNORECASE)
    data_match = re.search(r'data\s*=\s*(.+)$', text, flags=re.IGNORECASE)
    if type_match and data_match:
        chart_type = type_match.group(1).lower()
        data = _parse_chart_data_pairs(data_match.group(1))
        return (chart_type, data) if data else None
    return None

def _parse_chart_data_pairs(raw_data: str) -> dict[str, float]:
    entries = [entry.strip() for entry in re.split(r'[,;\n]+', str(raw_data or '')) if entry.strip()]
    parsed: dict[str, float] = {}
    seen: set[str] = set()
    for entry in entries:
        if ':' not in entry:
            return {}
        label, raw_value = entry.split(':', 1)
        key = label.strip()
        if not key:
            return {}
        normalized_key = key.casefold()
        if normalized_key in seen:
            return {}
        try:
            parsed[key] = float(raw_value.strip().replace('%', ''))
        except ValueError:
            return {}
        seen.add(normalized_key)
    return parsed

def _insert_docx_image(doc, image_path: Path, format_options: dict | None = None) -> None:
    try:
        if Inches is not None:
            doc.add_picture(str(image_path), width=Inches(6))
        else:
            doc.add_picture(str(image_path))
        if doc.paragraphs:
            _apply_docx_formatting(doc.paragraphs[-1], format_options)
    except Exception:
        logger.warning('Failed to insert image into docx: %s', image_path)

def _parse_hyperlink_detail(value: str) -> tuple[str, str]:
    raw = str(value or '').strip()
    if not raw:
        return '', ''
    if '|' in raw:
        text, url = raw.split('|', 1)
        text = text.strip()
        url = url.strip()
        if _is_http_url(url):
            return text, url
    if _is_http_url(raw):
        return raw, raw
    for token in raw.split():
        if _is_http_url(token):
            label = raw.replace(token, '').strip()
            return label or token, token
    return '', ''

def _is_http_url(value: str) -> bool:
    lowered = str(value or '').strip().lower()
    return lowered.startswith('http://') or lowered.startswith('https://')

def _add_docx_hyperlink(paragraph, url: str, text: str) -> None:
    if OxmlElement is None or qn is None or RELATIONSHIP_TYPE is None:
        paragraph.add_run(f'{text} ({url})')
        return
    try:
        relationship_id = paragraph.part.relate_to(url, RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
        hyperlink = OxmlElement('w:hyperlink')
        hyperlink.set(qn('r:id'), relationship_id)

        run_element = OxmlElement('w:r')
        run_properties = OxmlElement('w:rPr')

        color = OxmlElement('w:color')
        color.set(qn('w:val'), '0563C1')
        run_properties.append(color)

        underline = OxmlElement('w:u')
        underline.set(qn('w:val'), 'single')
        run_properties.append(underline)

        run_element.append(run_properties)
        text_element = OxmlElement('w:t')
        text_element.text = text
        run_element.append(text_element)
        hyperlink.append(run_element)

        paragraph._p.append(hyperlink)
    except Exception:
        paragraph.add_run(f'{text} ({url})')

def _set_docx_header_text(doc, value: str, format_options: dict | None = None) -> None:
    for section in doc.sections:
        header = section.header
        paragraph = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        if paragraph.text.strip():
            paragraph = header.add_paragraph()
        paragraph.text = str(value)
        _apply_docx_formatting(paragraph, format_options)

def _set_docx_footer_text(doc, value: str, format_options: dict | None = None) -> None:
    for section in doc.sections:
        footer = section.footer
        paragraph = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        if paragraph.text.strip():
            paragraph = footer.add_paragraph()
        paragraph.text = str(value)
        _apply_docx_formatting(paragraph, format_options)

def _add_docx_page_numbers(doc, format_options: dict | None = None) -> None:
    for section in doc.sections:
        footer = section.footer
        paragraph = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        if paragraph.text.strip():
            paragraph = footer.add_paragraph()
        paragraph.add_run('Page ')
        _append_field_code(paragraph.add_run(), ' PAGE ')
        _apply_docx_formatting(paragraph, format_options)

def _add_docx_toc_field(doc, format_options: dict | None = None) -> None:
    heading = doc.add_paragraph('Table of Contents')
    _apply_docx_formatting(heading, format_options)
    paragraph = doc.add_paragraph()
    _append_field_code(paragraph.add_run(), ' TOC \\o "1-3" \\h \\z \\u ')
    _apply_docx_formatting(paragraph, format_options)

def _append_field_code(run, instruction: str) -> None:
    if OxmlElement is None or qn is None:
        run.text = instruction.strip()
        return
    begin = OxmlElement('w:fldChar')
    begin.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText')
    instr.set(qn('xml:space'), 'preserve')
    instr.text = instruction
    separate = OxmlElement('w:fldChar')
    separate.set(qn('w:fldCharType'), 'separate')
    end = OxmlElement('w:fldChar')
    end.set(qn('w:fldCharType'), 'end')
    run._r.append(begin)
    run._r.append(instr)
    run._r.append(separate)
    run._r.append(end)

def _apply_docx_margins(doc, value: str) -> None:
    parsed = _parse_box_lengths(value)
    if not parsed:
        return
    for section in doc.sections:
        if parsed.get('top') is not None:
            section.top_margin = parsed['top']
        if parsed.get('right') is not None:
            section.right_margin = parsed['right']
        if parsed.get('bottom') is not None:
            section.bottom_margin = parsed['bottom']
        if parsed.get('left') is not None:
            section.left_margin = parsed['left']

def _apply_docx_paddings(doc, value: str) -> None:
    parsed_pt = _parse_box_points(value, default=8.0)
    if not parsed_pt:
        return
    for paragraph in doc.paragraphs:
        format_obj = paragraph.paragraph_format
        if Pt is not None:
            format_obj.space_before = Pt(parsed_pt['top'])
            format_obj.space_after = Pt(parsed_pt['bottom'])
            format_obj.left_indent = Pt(parsed_pt['left'])
            format_obj.right_indent = Pt(parsed_pt['right'])
    try:
        normal_style = doc.styles['Normal']
        if Pt is not None:
            normal_style.paragraph_format.space_before = Pt(parsed_pt['top'])
            normal_style.paragraph_format.space_after = Pt(parsed_pt['bottom'])
            normal_style.paragraph_format.left_indent = Pt(parsed_pt['left'])
            normal_style.paragraph_format.right_indent = Pt(parsed_pt['right'])
    except Exception:
        logger.warning('Failed to apply padding defaults to Normal style')

def _apply_docx_layout(doc, value: str) -> None:
    lowered = str(value or '').strip().lower()
    if not lowered:
        return
    for section in doc.sections:
        if WD_ORIENTATION is not None and 'landscape' in lowered:
            section.orientation = WD_ORIENTATION.LANDSCAPE
            section.page_width, section.page_height = section.page_height, section.page_width
        elif WD_ORIENTATION is not None and 'portrait' in lowered:
            section.orientation = WD_ORIENTATION.PORTRAIT
            section.page_width, section.page_height = section.page_height, section.page_width

        if qn is None:
            continue
        cols = section._sectPr.xpath('./w:cols')
        if cols:
            cols_element = cols[0]
        else:
            cols_element = OxmlElement('w:cols') if OxmlElement is not None else None
            if cols_element is not None:
                section._sectPr.append(cols_element)
        if cols_element is None:
            continue
        if '3 column' in lowered or 'three column' in lowered:
            cols_element.set(qn('w:num'), '3')
        elif '2 column' in lowered or 'two column' in lowered:
            cols_element.set(qn('w:num'), '2')
        elif 'single column' in lowered or '1 column' in lowered:
            cols_element.set(qn('w:num'), '1')

def _apply_docx_template(doc, value: str, format_options: dict | None = None) -> None:
    lowered = str(value or '').strip().lower()
    if not lowered:
        return
    if 'report' in lowered:
        _add_docx_toc_field(doc, format_options)
        _apply_docx_layout(doc, 'portrait, single column')
        return
    if 'resume' in lowered:
        _apply_docx_layout(doc, 'portrait, single column')
        _apply_docx_margins(doc, 'top:0.5in,right:0.5in,bottom:0.5in,left:0.5in')
        return
    if 'brochure' in lowered or 'flyer' in lowered:
        _apply_docx_layout(doc, 'landscape, 2 column')
        _apply_docx_margins(doc, 'top:0.4in,right:0.4in,bottom:0.4in,left:0.4in')
        return
    if 'article' in lowered:
        _apply_docx_layout(doc, 'portrait, single column')
        _apply_docx_margins(doc, 'top:1in,right:1in,bottom:1in,left:1in')
        return

def _apply_docx_theme(doc, value: str, format_options: dict | None = None) -> None:
    theme = str(value or '').strip().lower()
    if not theme:
        return
    presets = {
        'classic': {'font': 'Times New Roman', 'color': '#000000'},
        'modern': {'font': 'Calibri', 'color': '#1f2937'},
        'minimal': {'font': 'Arial', 'color': '#222222'},
        'bright': {'font': 'Verdana', 'color': '#0b5394'},
    }
    selected = presets.get(theme, None)
    if selected is None:
        selected = {
            'font': str((format_options or {}).get('font', 'Calibri')),
            'color': str((format_options or {}).get('color', '#1f2937')),
        }

    style_options = _normalize_format_options({
        'font': selected.get('font'),
        'color': selected.get('color'),
        'size': (format_options or {}).get('size'),
        'styles': (format_options or {}).get('styles', []),
        'alignment': (format_options or {}).get('alignment', ''),
    })
    for paragraph in doc.paragraphs:
        _apply_docx_formatting(paragraph, style_options)

def _apply_docx_background(doc, value: str) -> None:
    color = _parse_color_to_hex(value)
    if not color:
        return
    if OxmlElement is None or qn is None:
        return
    try:
        root = doc.part._element
        background = root.find(qn('w:background'))
        if background is None:
            background = OxmlElement('w:background')
            root.insert(0, background)
        background.set(qn('w:color'), color)
    except Exception:
        logger.warning('Failed to apply document background: %s', value)

def _apply_docx_borders(doc, value: str, format_options: dict | None = None) -> None:
    color = _parse_color_to_hex(value) or '333333'
    target = doc.add_paragraph('')
    _set_paragraph_border(target, color=color, size=8, style='single')
    _apply_docx_formatting(target, format_options)

def _set_paragraph_border(paragraph, color: str = '333333', size: int = 8, style: str = 'single') -> None:
    if OxmlElement is None or qn is None:
        return
    paragraph_properties = paragraph._p.get_or_add_pPr()
    border_properties = paragraph_properties.find(qn('w:pBdr'))
    if border_properties is None:
        border_properties = OxmlElement('w:pBdr')
        paragraph_properties.append(border_properties)
    for edge in ('top', 'left', 'bottom', 'right'):
        edge_element = border_properties.find(qn(f'w:{edge}'))
        if edge_element is None:
            edge_element = OxmlElement(f'w:{edge}')
            border_properties.append(edge_element)
        edge_element.set(qn('w:val'), style)
        edge_element.set(qn('w:sz'), str(max(size, 2)))
        edge_element.set(qn('w:space'), '1')
        edge_element.set(qn('w:color'), color)

def _add_docx_index_block(doc, value: str, format_options: dict | None = None) -> None:
    entries = sorted(set(_parse_detail_list(value)), key=str.casefold)
    heading = doc.add_paragraph('Index')
    _apply_docx_formatting(heading, format_options)
    for entry in entries:
        paragraph = doc.add_paragraph(entry)
        _apply_docx_formatting(paragraph, format_options)

def _add_docx_reference_block(
    doc,
    title: str,
    value: str,
    format_options: dict | None = None,
    numbered: bool = False,
) -> None:
    entries = _parse_detail_list(value)
    if not entries:
        return
    heading = doc.add_paragraph(title)
    _apply_docx_formatting(heading, format_options)
    for idx, entry in enumerate(entries, start=1):
        text = f'[{idx}] {entry}' if numbered else entry
        paragraph = doc.add_paragraph(text)
        _apply_docx_formatting(paragraph, format_options)

def _parse_detail_list(value: str) -> list[str]:
    raw = str(value or '').strip()
    if not raw:
        return []
    tokens = [piece.strip() for piece in re.split(r'[;\n/]+', raw) if piece.strip()]
    if len(tokens) <= 1 and ',' in raw:
        tokens = [piece.strip() for piece in raw.split(',') if piece.strip()]
    return tokens

def _parse_box_lengths(value: str) -> dict[str, object]:
    default_length = _parse_measurement(value)
    if default_length is not None:
        return {'top': default_length, 'right': default_length, 'bottom': default_length, 'left': default_length}

    sides: dict[str, object] = {}
    for part in re.split(r'[;,]+', str(value or '')):
        chunk = part.strip()
        if not chunk:
            continue
        if ':' in chunk:
            side, raw_measure = chunk.split(':', 1)
        elif '=' in chunk:
            side, raw_measure = chunk.split('=', 1)
        else:
            continue
        side_key = side.strip().lower()
        measurement = _parse_measurement(raw_measure.strip())
        if measurement is None:
            continue
        if side_key in ('top', 'right', 'bottom', 'left'):
            sides[side_key] = measurement
        elif side_key in ('all', 'margin', 'margins'):
            sides = {'top': measurement, 'right': measurement, 'bottom': measurement, 'left': measurement}
    return sides

def _parse_box_points(value: str, default: float = 8.0) -> dict[str, float]:
    single_value = _parse_numeric_value(value)
    if single_value is not None:
        return {'top': single_value, 'right': single_value, 'bottom': single_value, 'left': single_value}
    parsed = {'top': default, 'right': default, 'bottom': default, 'left': default}
    for part in re.split(r'[;,]+', str(value or '')):
        chunk = part.strip()
        if not chunk:
            continue
        if ':' in chunk:
            side, raw_value = chunk.split(':', 1)
        elif '=' in chunk:
            side, raw_value = chunk.split('=', 1)
        else:
            continue
        numeric = _parse_numeric_value(raw_value)
        if numeric is None:
            continue
        side_key = side.strip().lower()
        if side_key in parsed:
            parsed[side_key] = numeric
        elif side_key in ('all', 'padding', 'paddings'):
            parsed = {'top': numeric, 'right': numeric, 'bottom': numeric, 'left': numeric}
    return parsed

def _parse_numeric_value(value: str) -> float | None:
    match = re.search(r'(\d+(?:\.\d+)?)', str(value or ''))
    if not match:
        return None
    try:
        return float(match.group(1))
    except ValueError:
        return None

def _parse_measurement(value: str):
    text = str(value or '').strip().lower()
    if not text:
        return None
    match = re.search(r'(\d+(?:\.\d+)?)\s*(mm|cm|in|inch|inches|pt)?', text)
    if not match:
        return None
    amount = float(match.group(1))
    unit = match.group(2) or 'in'
    if unit in ('in', 'inch', 'inches'):
        return Inches(amount) if Inches is not None else None
    if unit == 'cm' and Cm is not None:
        return Cm(amount)
    if unit == 'mm' and Mm is not None:
        return Mm(amount)
    if unit == 'pt' and Pt is not None:
        return Pt(amount)
    return Inches(amount) if Inches is not None else None

def _parse_color_to_hex(value: str) -> str:
    rgb = _parse_color_to_rgb(value)
    if rgb is None:
        return ''
    return f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'

def _apply_docx_formatting(paragraph, format_options: dict | None = None) -> None:
    options = _normalize_format_options(format_options)
    if paragraph is None:
        return
    if WD_PARAGRAPH_ALIGNMENT is not None:
        alignment_map = {
            'left': WD_PARAGRAPH_ALIGNMENT.LEFT,
            'center': WD_PARAGRAPH_ALIGNMENT.CENTER,
            'right': WD_PARAGRAPH_ALIGNMENT.RIGHT,
            'justify': WD_PARAGRAPH_ALIGNMENT.JUSTIFY,
        }
        alignment = str(options.get('alignment', '')).strip().lower()
        if alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]

    color_rgb = _parse_color_to_rgb(options.get('color'))
    font_name = str(options.get('font', '')).strip()
    size = options.get('size')
    style_tokens = {str(item).lower() for item in options.get('styles', [])} if isinstance(options.get('styles'), list) else set()

    if not paragraph.runs:
        paragraph.add_run('')

    for run in paragraph.runs:
        if font_name:
            run.font.name = font_name
        if Pt is not None and size is not None:
            try:
                run.font.size = Pt(float(size))
            except Exception:
                pass
        if RGBColor is not None and color_rgb is not None:
            run.font.color.rgb = RGBColor(*color_rgb)
        if 'bold' in style_tokens:
            run.bold = True
        if 'italic' in style_tokens:
            run.italic = True
        if 'underline' in style_tokens:
            run.underline = True

def _parse_color_to_rgb(color_value: object) -> tuple[int, int, int] | None:
    if color_value is None:
        return None
    value = str(color_value).strip().lower()
    if not value:
        return None

    named = {
        'black': (0, 0, 0),
        'white': (255, 255, 255),
        'red': (255, 0, 0),
        'green': (0, 128, 0),
        'blue': (0, 0, 255),
        'yellow': (255, 255, 0),
        'orange': (255, 165, 0),
        'purple': (128, 0, 128),
        'gray': (128, 128, 128),
        'grey': (128, 128, 128),
        'brown': (165, 42, 42),
    }
    if value in named:
        return named[value]

    if re.fullmatch(r'#[0-9a-f]{3}', value):
        return (int(value[1] * 2, 16), int(value[2] * 2, 16), int(value[3] * 2, 16))

    if re.fullmatch(r'#[0-9a-f]{6}', value):
        return (int(value[1:3], 16), int(value[3:5], 16), int(value[5:7], 16))
    return None

def _parse_docx_blocks(text: str) -> list[dict[str, object]]:
    lines = text.splitlines()
    blocks: list[dict[str, object]] = []
    paragraph_buffer: list[str] = []
    index = 0

    while index < len(lines):
        current = lines[index]
        if not current.strip():
            _flush_paragraph_buffer(paragraph_buffer, blocks)
            index += 1
            continue

        table_candidate = _consume_table_block(lines, index)
        if table_candidate is not None:
            _flush_paragraph_buffer(paragraph_buffer, blocks)
            rows, next_index = table_candidate
            blocks.append({'type': 'table', 'rows': rows})
            logger.debug('Detected table block in input | start_line=%d rows=%d cols=%d', index, len(rows), len(rows[0]) if rows else 0)
            index = next_index
            continue

        paragraph_buffer.append(current)
        index += 1

    _flush_paragraph_buffer(paragraph_buffer, blocks)
    logger.debug('DOCX parse completed | input_lines=%d blocks=%d', len(lines), len(blocks))
    return blocks

def _flush_paragraph_buffer(paragraph_buffer: list[str], blocks: list[dict[str, object]]) -> None:
    for line in paragraph_buffer:
        value = line.strip()
        if value:
            blocks.append({'type': 'paragraph', 'value': value})
            logger.debug('Buffered paragraph added | preview=%s', preview_text(value))
    paragraph_buffer.clear()

def _consume_table_block(lines: list[str], start: int) -> tuple[list[list[str]], int] | None:
    for delimiter in ('|', '\t'):
        rows, next_index = _collect_table_rows(lines, start, delimiter)
        if not rows:
            continue
        if len(rows) < 2:
            continue
        if len(rows[0]) < 2:
            continue
        logger.debug(
            'Table delimiter selected | delimiter=%s start=%d rows=%d cols=%d',
            repr(delimiter),
            start,
            len(rows),
            len(rows[0]),
        )
        return rows, next_index
    return None

def _collect_table_rows(lines: list[str], start: int, delimiter: str) -> tuple[list[list[str]], int]:
    rows: list[list[str]] = []
    index = start

    while index < len(lines):
        raw_line = lines[index]
        if not raw_line.strip():
            break
        cells = _split_table_line(raw_line, delimiter)
        if len(cells) < 2:
            break
        rows.append(cells)
        index += 1

    if not rows:
        return [], start

    if len({len(row) for row in rows}) != 1:
        return [], start

    max_cols = max(len(row) for row in rows)
    normalized = [row + [''] * (max_cols - len(row)) for row in rows]

    if delimiter == '|' and len(normalized) > 1 and _is_markdown_separator_row(normalized[1]):
        normalized = [normalized[0]] + normalized[2:]

    if not normalized or len(normalized[0]) < 2:
        return [], start

    return normalized, index

def _split_table_line(line: str, delimiter: str) -> list[str]:
    candidate = line.strip()
    if delimiter == '|':
        if '|' not in candidate:
            return []
        if candidate.startswith('|'):
            candidate = candidate[1:]
        if candidate.endswith('|'):
            candidate = candidate[:-1]
        return [part.strip() for part in candidate.split('|')]
    if delimiter in (',', ';'):
        try:
            parsed = next(csv.reader([line], delimiter=delimiter))
        except Exception:
            return []
        return [part.strip() for part in parsed]
    return [part.strip() for part in line.split('\t')]

def _is_markdown_separator_row(row: list[str]) -> bool:
    if len(row) < 2:
        return False
    for cell in row:
        normalized = cell.replace(' ', '')
        if not normalized:
            return False
        if not re.fullmatch(r':?-{3,}:?', normalized):
            return False
    return True

def _read_docx_content(doc) -> str:
    parts: list[str] = []
    for block_type, block in _iter_docx_blocks(doc):
        if block_type == 'paragraph':
            if block.text.strip():
                parts.append(block.text)
            continue
        table = block
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells]
            if any(row_values):
                parts.append(' | '.join(row_values))
    logger.debug('Extracted docx content parts=%d', len(parts))
    return '\n'.join(parts)

def _iter_docx_blocks(doc):
    if CT_P is None or CT_Tbl is None:
        logger.warning('DOCX block iterator unavailable because python-docx XML classes are missing')
        return
    body = doc.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield 'paragraph', Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield 'table', Table(child, doc)

def _delete_text_content(file_name: str, delete_request: str) -> str:
    path = Path(file_name)
    if not path.exists():
        return f'Error deleting content: file not found: {file_name}'
    try:
        original = path.read_text(encoding='utf-8', errors='replace')
        updated, matches = _remove_text_matches(original, delete_request)
        logger.debug(
            'Text delete match scan | file=%s query_preview=%s matches=%d original_len=%d updated_len=%d',
            file_name,
            preview_text(delete_request),
            matches,
            len(original),
            len(updated),
        )
        if matches == 0:
            return f'No matching content found in "{file_name}".'
        normalized = _normalize_text_flow(updated)
        path.write_text(normalized, encoding='utf-8')
        logger.info('Deleted %d text matches from %s', matches, file_name)
        return f'Updated "{file_name}": removed {matches} matching text occurrence(s).'
    except Exception as exc:
        logger.exception('Failed to delete text content from %s', file_name)
        return f'Error deleting content: {exc}'

def _delete_docx_content(file_name: str, delete_request: str) -> str:
    if Document is None:
        msg = 'python-docx is required to update .docx files. Install with "pip install python-docx".'
        logger.error(msg)
        return 'Missing dependency: python-docx'

    path = Path(file_name)
    if not path.exists():
        return f'Error deleting content: file not found: {file_name}'

    try:
        doc = Document(file_name)
        table_mode, table_value = _parse_table_delete_request(delete_request)
        logger.info(
            'DOCX delete parsed table directive | file=%s mode=%s value=%s',
            file_name,
            table_mode,
            preview_text(table_value),
        )
        table_removed = _delete_docx_tables(doc, table_mode, table_value)

        text_targets = _extract_text_delete_targets(delete_request, table_mode)
        logger.info('DOCX delete parsed text targets | file=%s targets=%d', file_name, len(text_targets))
        text_removed = 0
        for target in text_targets:
            _, count = _remove_text_matches_from_docx(doc, target)
            text_removed += count
            logger.debug('DOCX text target processed | file=%s target=%s matches=%d', file_name, preview_text(target), count)

        empty_removed = _remove_empty_docx_paragraphs(doc)
        if table_removed == 0 and text_removed == 0 and empty_removed == 0:
            return f'No matching content found in "{file_name}".'

        doc.save(file_name)
        logger.info(
            'Updated docx %s by deleting %d table(s), %d text match(es), and %d empty paragraph(s)',
            file_name,
            table_removed,
            text_removed,
            empty_removed,
        )
        return (
            f'Updated "{file_name}": removed {table_removed} table(s), '
            f'{text_removed} text match(es), and {empty_removed} empty paragraph(s).'
        )
    except Exception as exc:
        logger.exception('Failed to delete docx content from %s', file_name)
        return f'Error deleting content: {exc}'

def _parse_table_delete_request(delete_request: str) -> tuple[str | None, str | int | None]:
    segments = [segment.strip() for segment in re.split(r'[;\n]+', delete_request) if segment.strip()]
    if not segments:
        return None, None

    prefixes = ('delete ', 'remove ', 'drop ')
    for segment in segments:
        cleaned = segment
        lowered = cleaned.lower()
        for prefix in prefixes:
            if lowered.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()
                lowered = cleaned.lower()
                break

        if not lowered.startswith('table'):
            continue

        remainder = cleaned[5:].strip()
        if remainder.startswith(':'):
            remainder = remainder[1:].strip()
        logger.debug('Evaluating table delete segment | segment=%s remainder=%s', preview_text(segment), preview_text(remainder))

        if not remainder:
            return 'all', None

        remainder_lower = remainder.lower()
        if remainder_lower in ('all', '*', 'tables'):
            return 'all', None
        if 'all' in remainder_lower and any(token in remainder_lower for token in ('table', 'content', 'rows')):
            return 'all', None
        if remainder_lower.startswith('and'):
            return 'all', None

        if remainder_lower.startswith('contains'):
            value = remainder[8:].strip()
            if value.startswith(':'):
                value = value[1:].strip()
            if value:
                return 'contains', value
            return 'all', None

        if remainder.isdigit():
            return 'index', int(remainder)

        match = re.search(r'\d+', remainder)
        if match and any(token in remainder_lower for token in ('index', '#', 'number', 'table')):
            return 'index', int(match.group(0))

        return 'contains', remainder

    return None, None

def _extract_text_delete_targets(delete_request: str, table_mode: str | None) -> list[str]:
    segments = [segment.strip() for segment in re.split(r'[;\n]+', delete_request) if segment.strip()]
    if not segments:
        return []

    targets: list[str] = []
    prefixes = ('delete ', 'remove ', 'drop ')
    for segment in segments:
        lowered = segment.lower()
        cleaned = segment
        for prefix in prefixes:
            if lowered.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()
                lowered = cleaned.lower()
                break
        if lowered.startswith('text:'):
            value = cleaned.split(':', 1)[1].strip()
            if value:
                targets.append(value)
            continue
        if table_mode is not None and lowered.startswith('table'):
            continue
        targets.append(cleaned)

    logger.debug('Extracted text delete targets | table_mode=%s targets=%s', table_mode, preview_text(targets))
    return targets

def _delete_docx_tables(doc, mode: str | None, value: str | int | None) -> int:
    if mode is None:
        return 0

    indices_to_remove: set[int] = set()
    tables = list(doc.tables)

    if mode == 'all':
        indices_to_remove.update(range(len(tables)))
    elif mode == 'index' and isinstance(value, int):
        target_index = value - 1
        if 0 <= target_index < len(tables):
            indices_to_remove.add(target_index)
    elif mode == 'contains' and isinstance(value, str):
        for idx, table in enumerate(tables):
            if _docx_table_contains(table, value):
                indices_to_remove.add(idx)

    for idx in sorted(indices_to_remove, reverse=True):
        table = tables[idx]
        table._element.getparent().remove(table._element)

    logger.info('DOCX table delete completed | mode=%s value=%s removed=%d', mode, preview_text(value), len(indices_to_remove))
    return len(indices_to_remove)

def _docx_table_contains(table, query: str) -> bool:
    query_lower = query.lower()
    for row in table.rows:
        for cell in row.cells:
            if query_lower in cell.text.lower():
                return True
    return False

def _remove_text_matches_from_docx(doc, query: str) -> tuple[bool, int]:
    changed = False
    total_matches = 0

    for paragraph in doc.paragraphs:
        matches = _remove_text_matches_from_docx_paragraph(paragraph, query)
        if matches > 0:
            changed = True
            total_matches += matches

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    matches = _remove_text_matches_from_docx_paragraph(paragraph, query)
                    if matches > 0:
                        changed = True
                        total_matches += matches

    logger.debug('DOCX text match removal completed | query=%s changed=%s matches=%d', preview_text(query), changed, total_matches)
    return changed, total_matches

def _remove_text_matches_from_docx_paragraph(paragraph, query: str) -> int:
    token = str(query or '').strip()
    if not token:
        return 0

    if not paragraph.runs:
        updated, matches = _remove_text_matches(paragraph.text, token)
        if matches > 0:
            paragraph.text = updated
            if len(paragraph.runs) == 1:
                paragraph.runs[0].text = _normalize_inline_spacing(paragraph.runs[0].text)
        return matches

    combined = ''.join(run.text for run in paragraph.runs)
    if not combined:
        return 0

    pattern = re.compile(re.escape(token), flags=re.IGNORECASE)
    matches = list(pattern.finditer(combined))
    if not matches:
        return 0

    ranges = _merge_ranges([(match.start(), match.end()) for match in matches])
    range_index = 0
    absolute_index = 0

    for run in paragraph.runs:
        if not run.text:
            continue
        kept_chars: list[str] = []
        for char in run.text:
            while range_index < len(ranges) and absolute_index >= ranges[range_index][1]:
                range_index += 1
            in_removed_range = (
                range_index < len(ranges)
                and ranges[range_index][0] <= absolute_index < ranges[range_index][1]
            )
            if not in_removed_range:
                kept_chars.append(char)
            absolute_index += 1
        run.text = ''.join(kept_chars)

    if len(paragraph.runs) == 1:
        paragraph.runs[0].text = _normalize_inline_spacing(paragraph.runs[0].text)

    return len(matches)

def _merge_ranges(ranges: list[tuple[int, int]]) -> list[tuple[int, int]]:
    if not ranges:
        return []
    ordered = sorted(ranges, key=lambda item: item[0])
    merged: list[tuple[int, int]] = [ordered[0]]
    for start, end in ordered[1:]:
        prev_start, prev_end = merged[-1]
        if start <= prev_end:
            merged[-1] = (prev_start, max(prev_end, end))
        else:
            merged.append((start, end))
    return merged

def _remove_text_matches(content: str, query: str) -> tuple[str, int]:
    token = query.strip()
    if not token:
        return content, 0
    pattern = re.compile(re.escape(token), flags=re.IGNORECASE)
    updated, matches = pattern.subn('', content)
    logger.debug('Regex substitution performed | query=%s matches=%d', preview_text(query), matches)
    return updated, matches

def _remove_empty_docx_paragraphs(doc) -> int:
    removed = 0
    for paragraph in list(doc.paragraphs):
        if paragraph.text.strip():
            continue
        parent = paragraph._element.getparent()
        if parent is None:
            continue
        parent.remove(paragraph._element)
        removed += 1
    logger.debug('Removed empty docx paragraphs | count=%d', removed)
    return removed

def _normalize_inline_spacing(text: str) -> str:
    collapsed = re.sub(r'[ \t]{2,}', ' ', text)
    collapsed = re.sub(r'\s+([,.;:!?])', r'\1', collapsed)
    return collapsed.strip()

def _normalize_text_flow(text: str) -> str:
    normalized = re.sub(r'[ \t]{2,}', ' ', text)
    normalized = re.sub(r'[ \t]+\n', '\n', normalized)
    normalized = re.sub(r'\n{3,}', '\n\n', normalized)
    logger.debug('Normalized text flow | original_len=%d normalized_len=%d', len(text), len(normalized))
    return normalized

###--BEGINNING OF CSV/PDF/PPT/MEDIA GENERATOR--###

def csv_write(file_name: str, text: str) -> str:
    path = Path(file_name)
    logger.info('CSV write requested | file=%s text_len=%d', file_name, len(str(text)))
    try:
        _ensure_parent_dir(path)
        rows = _parse_csv_rows(str(text))
        with open(file_name, 'w', encoding='utf-8', newline='') as handle:
            writer = csv.writer(handle)
            writer.writerows(rows)
        return file_name
    except Exception as exc:
        logger.exception('Failed to write csv file %s', file_name)
        return f'Error writing csv: {exc}'

def csv_append(file_name: str, text: str) -> str:
    path = Path(file_name)
    logger.info('CSV append requested | file=%s text_len=%d', file_name, len(str(text)))
    try:
        _ensure_parent_dir(path)
        rows = _parse_csv_rows(str(text))
        with open(file_name, 'a', encoding='utf-8', newline='') as handle:
            writer = csv.writer(handle)
            writer.writerows(rows)
        return file_name
    except Exception as exc:
        logger.exception('Failed to append csv file %s', file_name)
        return f'Error appending csv: {exc}'

def csv_read(file_name: str) -> str:
    logger.info('CSV read requested | file=%s', file_name)
    try:
        with open(file_name, 'r', encoding='utf-8', newline='') as handle:
            rows = list(csv.reader(handle))
        if not rows:
            return ''
        return '\n'.join(','.join(cell for cell in row) for row in rows)
    except Exception as exc:
        logger.exception('Failed to read csv file %s', file_name)
        return f'Error reading csv: {exc}'

def _parse_csv_rows(text: str) -> list[list[str]]:
    cleaned = str(text or '').strip()
    if not cleaned:
        return [['']]
    lines = [line for line in cleaned.splitlines() if line.strip()]
    if not lines:
        return [['']]

    candidate_delimiters = [',', '|', '\t', ';']
    best_delimiter = ','
    best_score = 1
    for delimiter in candidate_delimiters:
        score = max((len(_split_csv_like_line(line, delimiter)) for line in lines), default=1)
        if score > best_score:
            best_score = score
            best_delimiter = delimiter

    rows: list[list[str]] = []
    for line in lines:
        split = _split_csv_like_line(line, best_delimiter)
        if split:
            rows.append(split)

    return rows if rows else [[cleaned]]

def _split_csv_like_line(line: str, delimiter: str) -> list[str]:
    if delimiter in (',', ';'):
        try:
            parsed = next(csv.reader([line], delimiter=delimiter))
        except Exception:
            parsed = []
        return [cell.strip() for cell in parsed if cell is not None]
    if delimiter == '\t':
        return [cell.strip() for cell in line.split('\t')]
    if delimiter == '|':
        return [cell.strip() for cell in line.strip('|').split('|')]
    return [line.strip()]

def pdf_write(file_name: str, text: str) -> str:
    logger.info('PDF write requested | file=%s text_len=%d', file_name, len(str(text)))
    if reportlab_canvas is None:
        msg = 'Missing dependency: reportlab. Install with "pip install reportlab".'
        logger.error(msg)
        return msg
    path = Path(file_name)
    try:
        _ensure_parent_dir(path)
        _write_text_to_pdf(file_name, str(text))
        return file_name
    except Exception as exc:
        logger.exception('Failed to write pdf file %s', file_name)
        return f'Error writing pdf: {exc}'

def pdf_append(file_name: str, text: str) -> str:
    logger.info('PDF append requested | file=%s text_len=%d', file_name, len(str(text)))
    if reportlab_canvas is None:
        msg = 'Missing dependency: reportlab. Install with "pip install reportlab".'
        logger.error(msg)
        return msg
    existing_text = ''
    if Path(file_name).exists():
        previous = pdf_read(file_name)
        if isinstance(previous, str) and not previous.lower().startswith('error') and 'Missing dependency' not in previous:
            existing_text = previous
    combined = f'{existing_text}\n{text}'.strip()
    return pdf_write(file_name, combined)

def pdf_read(file_name: str) -> str:
    logger.info('PDF read requested | file=%s', file_name)
    if PdfReader is None:
        msg = 'Missing dependency: PyPDF2. Install with "pip install PyPDF2".'
        logger.error(msg)
        return msg
    try:
        reader = PdfReader(file_name)
        pages: list[str] = []
        for page in reader.pages:
            pages.append((page.extract_text() or '').strip())
        return '\n\n'.join(part for part in pages if part)
    except Exception as exc:
        logger.exception('Failed to read pdf file %s', file_name)
        return f'Error reading pdf: {exc}'

def _write_text_to_pdf(file_name: str, text: str) -> None:
    pdf = reportlab_canvas.Canvas(file_name)
    width, height = pdf._pagesize
    y = height - 50
    for raw_line in str(text or '').splitlines() or ['']:
        line = raw_line.rstrip()
        if y < 50:
            pdf.showPage()
            y = height - 50
        pdf.drawString(50, y, line[:1400])
        y -= 16
    pdf.save()

def ppt_write(file_name: str, text: str) -> str:
    logger.info('PPT write requested | file=%s text_len=%d', file_name, len(str(text)))
    if Presentation is None:
        msg = 'Missing dependency: python-pptx. Install with "pip install python-pptx".'
        logger.error(msg)
        return msg
    target = _normalize_ppt_target(file_name)
    try:
        _ensure_parent_dir(target)
        presentation = Presentation()
        _add_text_to_presentation(presentation, str(text))
        presentation.save(str(target))
        return str(target)
    except Exception as exc:
        logger.exception('Failed to write ppt file %s', file_name)
        return f'Error writing ppt: {exc}'

def ppt_append(file_name: str, text: str) -> str:
    logger.info('PPT append requested | file=%s text_len=%d', file_name, len(str(text)))
    if Presentation is None:
        msg = 'Missing dependency: python-pptx. Install with "pip install python-pptx".'
        logger.error(msg)
        return msg
    target = _normalize_ppt_target(file_name)
    try:
        _ensure_parent_dir(target)
        presentation = Presentation(str(target)) if target.exists() else Presentation()
        _add_text_to_presentation(presentation, str(text))
        presentation.save(str(target))
        return str(target)
    except Exception as exc:
        logger.exception('Failed to append ppt file %s', file_name)
        return f'Error appending ppt: {exc}'

def ppt_read(file_name: str) -> str:
    logger.info('PPT read requested | file=%s', file_name)
    if Presentation is None:
        msg = 'Missing dependency: python-pptx. Install with "pip install python-pptx".'
        logger.error(msg)
        return msg
    target = _normalize_ppt_target(file_name)
    try:
        presentation = Presentation(str(target))
        lines: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines.append(f'Slide {slide_index}:')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and str(shape.text).strip():
                    lines.append(shape.text.strip())
        return '\n'.join(lines).strip()
    except Exception as exc:
        logger.exception('Failed to read ppt file %s', file_name)
        return f'Error reading ppt: {exc}'

def _normalize_ppt_target(file_name: str) -> Path:
    path = Path(file_name)
    if path.suffix.lower() == '.ppt':
        return path.with_suffix('.pptx')
    if path.suffix.lower() != '.pptx':
        return path.with_suffix('.pptx')
    return path

def _add_text_to_presentation(presentation, text: str) -> None:
    chunks = [chunk.strip() for chunk in re.split(r'\n{2,}', str(text or '')) if chunk.strip()]
    if not chunks:
        chunks = ['']
    for index, chunk in enumerate(chunks):
        lines = [line.strip() for line in chunk.splitlines() if line.strip()]
        title = lines[0][:120] if lines else f'Slide {len(presentation.slides) + 1}'
        body_lines = lines[1:] if len(lines) > 1 else []
        layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
            slide.shapes.title.text = title
        body_placeholder = None
        if len(slide.placeholders) > 1:
            try:
                body_placeholder = slide.placeholders[1]
            except Exception:
                body_placeholder = None
        if body_placeholder is not None and hasattr(body_placeholder, 'text_frame'):
            text_frame = body_placeholder.text_frame
            text_frame.clear()
            if body_lines:
                text_frame.text = body_lines[0]
                for line in body_lines[1:]:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = line
            else:
                text_frame.text = ''
        elif body_lines:
            # Fall back to a simple newline body when no placeholder is available.
            if hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
                slide.shapes.title.text = f'{title}\n' + '\n'.join(body_lines[:6])

def media_read(file_name: str, media_type: str = 'media') -> str:
    path = Path(file_name)
    if not path.exists():
        return f'Error reading {media_type}: file not found: {file_name}'
    try:
        size = path.stat().st_size
        ext = path.suffix.lower().lstrip('.')
        details = [
            f'Media type: {media_type}',
            f'File: {path.resolve()}',
            f'Format: .{ext}' if ext else 'Format: unknown',
            f'Size: {size} bytes',
        ]
        if media_type == 'audio' and path.suffix.lower() == '.wav':
            with wave.open(str(path), 'rb') as handle:
                frames = handle.getnframes()
                framerate = handle.getframerate() or 1
                duration = frames / float(framerate)
                details.append(f'Duration: {duration:.2f} seconds')
                details.append(f'Channels: {handle.getnchannels()}')
                details.append(f'Sample rate: {framerate} Hz')
        return '\n'.join(details)
    except Exception as exc:
        logger.exception('Failed to read %s media file %s', media_type, file_name)
        return f'Error reading {media_type}: {exc}'

def media_write(file_name: str, text: str, media_type: str = 'media', mode: str = 'write') -> str:
    source = _extract_source_path(text)
    if not source:
        return (
            f'Error writing {media_type}: provide source path using '
            f'"path:/full/source/file" or "source:/full/source/file".'
        )
    source_path = Path(source).expanduser()
    if not source_path.is_absolute():
        source_path = (Path.cwd() / source_path).resolve()
    if not source_path.exists():
        return f'Error writing {media_type}: source file not found: {source_path}'

    target_path = Path(file_name)
    try:
        _ensure_parent_dir(target_path)
        if mode == 'append':
            return f'Error appending {media_type}: append is not supported; use write with source path.'
        shutil.copyfile(str(source_path), str(target_path))
        return str(target_path)
    except Exception as exc:
        logger.exception('Failed to write %s media file %s from %s', media_type, file_name, source_path)
        return f'Error writing {media_type}: {exc}'

def _extract_source_path(text: str) -> str:
    raw = str(text or '').strip()
    if not raw:
        return ''
    lower = raw.lower()
    for prefix in ('path:', 'source:'):
        if lower.startswith(prefix):
            return raw[len(prefix):].strip()
    return ''

###--END OF CSV/PDF/PPT/MEDIA GENERATOR--###

###--END OF TEXT FILE GENERATOR--###
###--BEGINNING OF EXCEL FILE GENERATOR --###

def xlsx_write(file_name: str, text) -> str:
    logger.info('XLSX write requested | file=%s text_len=%d preview=%s', file_name, len(str(text)), preview_text(text))
    if pd is None:
        msg = 'pandas is required to generate .xlsx files. Install with "pip install pandas openpyxl".'
        logger.error(msg)
        return 'Missing dependency: pandas'

    path = Path(file_name)
    try:
        _ensure_parent_dir(path)
        df = _text_to_dataframe(str(text))
        df.to_excel(file_name, index=False)
        logger.info('Created xlsx file %s', file_name)
        return file_name
    except Exception as exc:
        logger.exception('Failed to create xlsx file %s', file_name)
        return f'Error creating xlsx: {exc}'

def xlsx_read(file_name: str) -> str:
    ## read and summarize the content of the file using openai API
    logger.info('XLSX read requested | file=%s', file_name)
    if pd is None:
        msg = 'pandas is required to read .xlsx files. Install with "pip install pandas openpyxl".'
        logger.error(msg)
        return 'Missing dependency: pandas'
    try:
        df = pd.read_excel(file_name)
        output = df.to_string(index=False)
        logger.info('Read xlsx file %s', file_name)
        return output
    except Exception as exc:
        logger.exception('Failed to read xlsx file %s', file_name)
        return f'Error reading xlsx: {exc}'

def xlsx_append(file_name: str, text) -> str:
    logger.info('XLSX append requested | file=%s text_len=%d preview=%s', file_name, len(str(text)), preview_text(text))
    if pd is None:
        msg = 'pandas is required to append to .xlsx files. Install with "pip install pandas openpyxl".'
        logger.error(msg)
        return 'Missing dependency: pandas'
    path = Path(file_name)
    try:
        _ensure_parent_dir(path)
        try:
            existing = pd.read_excel(file_name)
        except FileNotFoundError:
            existing = pd.DataFrame(columns=['content'])

        new_df = _text_to_dataframe(str(text))
        all_columns = list(dict.fromkeys(list(existing.columns) + list(new_df.columns)))
        existing = existing.reindex(columns=all_columns, fill_value='')
        new_df = new_df.reindex(columns=all_columns, fill_value='')
        updated = pd.concat([existing, new_df], ignore_index=True)
        updated.to_excel(file_name, index=False)
        logger.info('Appended to xlsx file %s', file_name)
        return file_name
    except Exception as exc:
        logger.exception('Failed to append to xlsx file %s', file_name)
        return f'Error appending to xlsx: {exc}'

def _ensure_parent_dir(path: Path) -> None:
    # Creates any missing parent directories for nested paths
    if path.parent and not path.parent.exists():
        path.parent.mkdir(parents=True, exist_ok=True)
        logger.debug('Created parent directory for path %s', path.parent)

def _parse_cloud_reference(file_name: str) -> dict[str, str] | None:
    raw = str(file_name or '').strip()
    if not raw or ':' not in raw:
        return None
    prefix, remainder = raw.split(':', 1)
    provider_key = CLOUD_PROVIDER_ALIASES.get(prefix.strip().lower())
    if not provider_key:
        return None
    relative = remainder.strip().lstrip('/\\')
    if not relative:
        return None
    return {'provider': provider_key, 'relative': relative}

def _cloud_ref_with_result_suffix(cloud_ref: dict[str, str], result_path: Path) -> dict[str, str]:
    relative_path = Path(cloud_ref['relative'])
    result_suffix = result_path.suffix
    if not result_suffix:
        return dict(cloud_ref)
    if relative_path.suffix.lower() == result_suffix.lower():
        return dict(cloud_ref)
    if relative_path.suffix:
        updated_relative = str(relative_path.with_suffix(result_suffix))
    else:
        updated_relative = str(relative_path.parent / f'{relative_path.name}{result_suffix}')
    return {'provider': cloud_ref['provider'], 'relative': updated_relative}

def _cloud_display_name(original_name: str, cloud_ref: dict[str, str]) -> str:
    if ':' not in str(original_name):
        return original_name
    prefix = str(original_name).split(':', 1)[0].strip()
    return f'{prefix}:{cloud_ref["relative"]}'

def _rclone_is_available() -> bool:
    global _RCLONE_AVAILABLE
    if _RCLONE_AVAILABLE is not None:
        return _RCLONE_AVAILABLE
    try:
        completed = subprocess.run(
            ['rclone', 'version'],
            check=False,
            capture_output=True,
            text=True,
        )
        _RCLONE_AVAILABLE = completed.returncode == 0
    except Exception:
        _RCLONE_AVAILABLE = False
    if not _RCLONE_AVAILABLE:
        logger.info('rclone not available; cloud-prefixed paths use local/mounted folders only')
    return _RCLONE_AVAILABLE

def _rclone_remote_path(cloud_ref: dict[str, str]) -> str:
    provider = cloud_ref['provider']
    relative = cloud_ref['relative']
    env_name = f'FILEGEN_CLOUD_REMOTE_{provider.upper()}'
    remote_name = os.getenv(env_name) or RCLONE_REMOTE_DEFAULTS.get(provider, provider)
    return f'{remote_name}:{relative}'

def _rclone_download_remote(cloud_ref: dict[str, str], local_file: str) -> tuple[bool, str]:
    remote_path = _rclone_remote_path(cloud_ref)
    local_path = Path(local_file)
    _ensure_parent_dir(local_path)
    try:
        completed = subprocess.run(
            ['rclone', 'copyto', remote_path, str(local_path)],
            check=False,
            capture_output=True,
            text=True,
        )
        if completed.returncode == 0:
            return True, ''
        err = (completed.stderr or completed.stdout or '').strip()
        return False, err or f'Failed to download {remote_path}'
    except Exception as exc:
        return False, str(exc)

def _rclone_upload_remote(cloud_ref: dict[str, str], local_file: str) -> tuple[bool, str]:
    remote_path = _rclone_remote_path(cloud_ref)
    local_path = Path(local_file)
    if not local_path.exists():
        return False, f'Local file not found for upload: {local_file}'
    try:
        completed = subprocess.run(
            ['rclone', 'copyto', str(local_path), remote_path],
            check=False,
            capture_output=True,
            text=True,
        )
        if completed.returncode == 0:
            return True, ''
        err = (completed.stderr or completed.stdout or '').strip()
        return False, err or f'Failed to upload {remote_path}'
    except Exception as exc:
        return False, str(exc)

def _rclone_delete_remote(cloud_ref: dict[str, str]) -> tuple[bool, str]:
    remote_path = _rclone_remote_path(cloud_ref)
    try:
        completed = subprocess.run(
            ['rclone', 'deletefile', remote_path],
            check=False,
            capture_output=True,
            text=True,
        )
        if completed.returncode == 0:
            return True, ''
        err = (completed.stderr or completed.stdout or '').strip()
        return False, err or f'Failed to delete {remote_path}'
    except Exception as exc:
        return False, str(exc)

def _is_full_delete_request(delete_request: str) -> bool:
    request = str(delete_request or '').strip().lower()
    aliases = {'', 'file', 'entire file', 'full file', 'whole file', 'delete file', 'remove file'}
    return request in aliases

def _is_error_result(result: object) -> bool:
    if not isinstance(result, str):
        return False
    lower = result.lower()
    return (
        lower.startswith('error')
        or 'missing dependency' in lower
        or 'unsupported file type' in lower
        or 'invalid action' in lower
        or 'file not found' in lower
    )

def _resolve_storage_target(file_name: str) -> str:
    raw = str(file_name or '').strip()
    if not raw:
        return raw
    if ':' not in raw:
        return raw
    prefix, remainder = raw.split(':', 1)
    provider_key = CLOUD_PROVIDER_ALIASES.get(prefix.strip().lower())
    if not provider_key:
        return raw
    relative = remainder.strip().lstrip('/\\')
    if not relative:
        return raw
    base = _provider_storage_base(provider_key)
    target = base / relative
    _ensure_parent_dir(target)
    return str(target)

def _provider_storage_base(provider_key: str) -> Path:
    home = Path.home()
    candidates_by_provider = {
        'google_drive': [home / 'Google Drive', home / 'My Drive'],
        'dropbox': [home / 'Dropbox'],
        'onedrive': [home / 'OneDrive'],
        'icloud': [home / 'Library' / 'Mobile Documents' / 'com~apple~CloudDocs'],
        'box': [home / 'Box'],
        'pcloud': [home / 'pCloud Drive'],
        'mega': [home / 'MEGA'],
        'sync': [home / 'Sync'],
        'nextcloud': [home / 'Nextcloud'],
        'owncloud': [home / 'ownCloud'],
        'tresorit': [home / 'Tresorit'],
        's3': [Path.cwd() / 'cloud_storage' / 's3'],
    }
    candidates = candidates_by_provider.get(provider_key, [])
    for candidate in candidates:
        if candidate.exists():
            return candidate
    fallback = Path.cwd() / 'cloud_storage' / provider_key
    fallback.mkdir(parents=True, exist_ok=True)
    return fallback

def _text_to_dataframe(text: str):
    # Converts free-form text into a DataFrame, supporting multi-column tabular content.
    lines = [line for line in text.splitlines() if line.strip()]
    logger.debug('Converting text to dataframe | non_empty_lines=%d', len(lines))
    if not lines:
        return pd.DataFrame({'content': ['']})

    # Detect delimiter that yields the most columns
    delimiters = ['\t', ',', '|', ';']
    best_delim = None
    best_cols = 1
    for delim in delimiters:
        cols = max(len(line.split(delim)) for line in lines)
        if cols > best_cols:
            best_cols = cols
            best_delim = delim

    if best_delim is None or best_cols == 1:
        return pd.DataFrame({'content': lines})

    split_rows = [line.split(best_delim) for line in lines]
    max_len = max(len(r) for r in split_rows)
    padded_rows = [r + [''] * (max_len - len(r)) for r in split_rows]

    header_candidate = padded_rows[0]
    header_is_valid = len(set(header_candidate)) == len(header_candidate) and all(cell.strip() for cell in header_candidate)
    if header_is_valid:
        columns = header_candidate
        data_rows = padded_rows[1:] if len(padded_rows) > 1 else []
    else:
        columns = [f'Column{i+1}' for i in range(max_len)]
        data_rows = padded_rows

    df = pd.DataFrame(data_rows, columns=columns)
    logger.info('Parsed tabular content with %d columns using delimiter "%s"', len(columns), repr(best_delim))
    return df

###--END OF EXCEL FILE GENERATOR --###
##--BEGINNING OF IMAGE FILE GENERATOR --###

def image_creation(user_input: str, file_name: str, action:str) -> str:
   logger.info('Image creation requested | file=%s action=%s prompt_preview=%s', file_name, action, preview_text(user_input))
   return generate_image(user_input, file_name, action)

def image_read(user_input: str, file_name: str, action: str) -> str:
    logger.info('Image read requested | file=%s action=%s request_preview=%s', file_name, action, preview_text(user_input))
    return generate_image(user_input, file_name, action)


def image_append(user_input: str, file_name: str, action: str) -> str:
    # I wonder if we can append to an image by generating a new one based on the old one + new prompt, then saving it with the same name to overwrite it?
    # the answer is yes, we can do that. we can use the old image as a reference for the new one by including it in the prompt, and then save the new image with the same name to overwrite it.
    # what if the user doent like the new image?
    # if the user doesn't like the new image, they can simply generate a new one with a different prompt and save it with a different name. The original image will still be there if they want to keep it.
    logger.info('Image append/edit requested | file=%s action=%s prompt_preview=%s', file_name, action, preview_text(user_input))
    return generate_image(user_input, file_name, action)

##--END OF IMAGE FILE GENERATOR --###
##implemented font to docx file.
# working on it hold on

def _resolve_chart_output_path(file_name: str) -> Path:
    path = Path(file_name)
    if path.suffix.lower() not in ('.png', '.jpg', '.jpeg', '.webp'):
        path = path.with_suffix('.png')
    _ensure_parent_dir(path)
    return path

# Function to generate and save a chart
def generate_chart(chart_type: str, data: dict, output_path: str):
    """
    Generates a chart and saves it as an image.

    Args:
        chart_type (str): Type of chart (e.g., 'bar', 'pie').
        data (dict): Data for the chart (keys as labels, values as data points).
        output_path (str): Path to save the generated chart image.
    """
    try:
        if plt is None:
            raise RuntimeError('matplotlib is required for chart generation. Install with "pip install matplotlib".')
        chart_type = (chart_type or '').strip().lower()
        labels = [str(k) for k in data.keys()]
        values = [float(v) for v in data.values()]
        if not labels or not values:
            raise ValueError('Chart data cannot be empty.')

        if chart_type == 'bar':
            plt.figure(figsize=(10, 6))
            plt.bar(labels, values)
        elif chart_type == 'line':
            plt.figure(figsize=(10, 6))
            plt.plot(labels, values, marker='o')
        elif chart_type == 'scatter':
            plt.figure(figsize=(10, 6))
            plt.scatter(labels, values)
        elif chart_type == 'pie':
            if any(value < 0 for value in values):
                raise ValueError('Pie chart values cannot be negative.')
            if sum(values) <= 0:
                raise ValueError('Pie chart requires a total greater than zero.')
            plt.figure(figsize=(8, 8))
            plt.pie(values, labels=labels, autopct='%1.1f%%')  # Convert keys to list
        else:
            raise ValueError(f"Unsupported chart type: {chart_type}")

        plt.title(f"{chart_type.capitalize()} Chart")
        if chart_type != 'pie':
            plt.xlabel('Labels')
            plt.ylabel('Values')
            plt.xticks(rotation=45)
            plt.tight_layout()
        plt.savefig(output_path)
        plt.close()
        logger.info(f"Chart saved at {output_path}")
    except Exception as e:
        logger.error(f"Failed to generate chart: {e}")
        raise
